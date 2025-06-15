from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
import os
import random
import string
import tkinter as tk
from tkinter import filedialog, ttk, colorchooser, messagebox
import winreg
from reportlab.lib.pagesizes import letter
import shutil
from pptx.util import Pt
# Helpers for font and pptx picking, plus PDF writing

def get_installed_fonts():
    fonts = set()
    font_keys = [
        r"SOFTWARE\Microsoft\Windows NT\CurrentVersion\Fonts",
        r"SOFTWARE\Microsoft\Windows\CurrentVersion\Fonts"
    ]
    for key in font_keys:
        try:
            with winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE, key) as reg_key:
                i = 0
                while True:
                    try:
                        font_name, font_file, _ = winreg.EnumValue(reg_key, i)
                        fonts.add((font_name.split(" (")[0], font_file))
                        i += 1
                    except OSError:
                        break
        except FileNotFoundError:
            continue
    return sorted(fonts)

def choose_font_dialog(parent, fonts):
    font_choice = tk.StringVar(value=fonts[0][0])
    dlg = tk.Toplevel(parent)
    dlg.title("Choose a Font")
    dlg.geometry("420x160")
    dlg.resizable(False, False)

    tk.Label(dlg, text="Select a font:").pack(pady=(10, 0))
    dropdown = ttk.Combobox(dlg, textvariable=font_choice,
                            values=[f[0] for f in fonts] + ["Other..."],
                            state="readonly", width=40)
    dropdown.pack(pady=10)
    dropdown.focus_set()

    selected_path = {"path": None}

    def on_select():
        choice = font_choice.get()
        if choice == "Other...":
            dlg.withdraw()
            file_path = filedialog.askopenfilename(
                title="Select font file",
                filetypes=[("Font files", "*.ttf *.otf"), ("All files", "*.*")]
            )
            dlg.deiconify()
            if file_path:
                selected_path["path"] = file_path
                dlg.destroy()
        else:
            for name, path in fonts:
                if name == choice:
                    selected_path["path"] = os.path.join("C:\\Windows\\Fonts", path)
                    break
            dlg.destroy()

    tk.Button(dlg, text="Select Font", command=on_select).pack(pady=(0, 10))
    dlg.protocol("WM_DELETE_WINDOW", dlg.destroy)
    parent.wait_window(dlg)
    return selected_path["path"]

def register_custom_font(font_path):
    font_name = os.path.splitext(os.path.basename(font_path))[0]  # "Xaihand-Regular"
    if font_name not in pdfmetrics.getRegisteredFontNames():
        pdfmetrics.registerFont(TTFont(font_name, font_path))
    return font_name

def is_small_text_shape(shape, min_width=Pt(30), min_height=Pt(15)):
    # Filter out shapes with very small size, probably part of table cells
    shape_width = shape.width
    shape_height = shape.height
    if shape_width < min_width or shape_height < min_height:
        return True
    return False

def is_small_text_shape(shape, min_width=Pt(30), min_height=Pt(15)):
    shape_width = shape.width
    shape_height = shape.height
    if shape_width < min_width or shape_height < min_height:
        return True
    return False

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    notes = []
    for slide in prs.slides:
        slide_text = ""
        text_shapes = [shape for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        
        for shape in text_shapes:
            text = shape.text.strip()
            if len(text) <= 2:  # remove text sections with 2 or fewer chars
                continue
            if is_small_text_shape(shape):
                continue
            slide_text += text + "\n"
        notes.append(slide_text.strip())
    return notes

def write_notes_to_pdf(notes, output_path, font_path, font_name, font_size, margins, line_spacing, font_color_hex):
    from reportlab.lib.colors import HexColor
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    import os

    def register_custom_font(path):
        name = os.path.splitext(os.path.basename(path))[0]
        if name not in pdfmetrics.getRegisteredFontNames():
            pdfmetrics.registerFont(TTFont(name, path))
        return name

    # Register and use the actual internal font name
    actual_font_name = register_custom_font(font_path)

    c = canvas.Canvas(output_path, pagesize=letter)
    c.setFont(actual_font_name, font_size)
    c.setFillColor(HexColor(font_color_hex))

    width, height = letter
    x = margins['left']
    y = height - margins['top']
    bottom_margin = margins['bottom']

    for slide_text in notes:
        for line in slide_text.split('\n'):
            if y < bottom_margin:
                c.showPage()
                c.setFont(actual_font_name, font_size)
                c.setFillColor(HexColor(font_color_hex))
                y = height - margins['top']
            c.drawString(x, y, line)
            y -= line_spacing

        # Add a soft space between slides, just like in your old code 💋
        y -= line_spacing

    c.save()

# Main UI class

class NoteGenUI(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("PowerPoint Notes to PDF Generator")
        self.geometry("600x550")
        self.resizable(False, False)

        # Variables
        self.pptx_path = tk.StringVar()
        self.font_path = None
        self.font_name = tk.StringVar(value="Select a font")
        self.font_size = tk.IntVar(value=12)
        self.line_spacing = tk.DoubleVar(value=14.0)
        self.font_color = (0, 0, 0)  # black default
        self.margin_left = tk.IntVar(value=72)
        self.margin_right = tk.IntVar(value=72)
        self.margin_top = tk.IntVar(value=72)
        self.margin_bottom = tk.IntVar(value=72)
        self.output_pdf = tk.StringVar()

        self.installed_fonts = get_installed_fonts()

        self.create_widgets()

    def create_widgets(self):
        pad = 10
        # PPTX file selection
        tk.Label(self, text="PowerPoint file:").pack(anchor="w", padx=pad, pady=(pad, 0))
        frame_pptx = tk.Frame(self)
        frame_pptx.pack(fill="x", padx=pad)
        tk.Entry(frame_pptx, textvariable=self.pptx_path, state="readonly").pack(side="left", fill="x", expand=True)
        tk.Button(frame_pptx, text="Browse", command=self.browse_pptx).pack(side="left", padx=5)

        # Font selection
        tk.Label(self, text="Font:").pack(anchor="w", padx=pad, pady=(pad, 0))
        frame_font = tk.Frame(self)
        frame_font.pack(fill="x", padx=pad)
        tk.Label(frame_font, textvariable=self.font_name).pack(side="left")
        tk.Button(frame_font, text="Choose Font", command=self.choose_font).pack(side="left", padx=5)

        # Font size
        tk.Label(self, text="Font Size:").pack(anchor="w", padx=pad, pady=(pad, 0))
        tk.Spinbox(self, from_=6, to=72, textvariable=self.font_size, width=5).pack(anchor="w", padx=pad)

        # Font color
        frame_color = tk.Frame(self)
        frame_color.pack(anchor="w", padx=pad, pady=(pad, 0))
        tk.Label(frame_color, text="Font Color:").pack(side="left")
        self.color_display = tk.Canvas(frame_color, width=40, height=20, bg=self._color_to_hex(self.font_color), bd=1, relief="sunken")
        self.color_display.pack(side="left", padx=5)
        tk.Button(frame_color, text="Choose Color", command=self.choose_color).pack(side="left")

        # Line spacing
        tk.Label(self, text="Line Spacing:").pack(anchor="w", padx=pad, pady=(pad, 0))
        tk.Spinbox(self, from_=10.0, to=40.0, increment=0.5, textvariable=self.line_spacing, width=5).pack(anchor="w", padx=pad)

        # Margins
        tk.Label(self, text="Margins (points):").pack(anchor="w", padx=pad, pady=(pad, 0))
        frame_margins = tk.Frame(self)
        frame_margins.pack(anchor="w", padx=pad)

        tk.Label(frame_margins, text="Left:").grid(row=0, column=0, padx=3, pady=2)
        tk.Spinbox(frame_margins, from_=0, to=200, textvariable=self.margin_left, width=5).grid(row=0, column=1, padx=3, pady=2)

        tk.Label(frame_margins, text="Right:").grid(row=0, column=2, padx=3, pady=2)
        tk.Spinbox(frame_margins, from_=0, to=200, textvariable=self.margin_right, width=5).grid(row=0, column=3, padx=3, pady=2)

        tk.Label(frame_margins, text="Top:").grid(row=1, column=0, padx=3, pady=2)
        tk.Spinbox(frame_margins, from_=0, to=200, textvariable=self.margin_top, width=5).grid(row=1, column=1, padx=3, pady=2)

        tk.Label(frame_margins, text="Bottom:").grid(row=1, column=2, padx=3, pady=2)
        tk.Spinbox(frame_margins, from_=0, to=200, textvariable=self.margin_bottom, width=5).grid(row=1, column=3, padx=3, pady=2)

        # Output PDF selection
        tk.Label(self, text="Output PDF file:").pack(anchor="w", padx=pad, pady=(pad, 0))
        frame_output = tk.Frame(self)
        frame_output.pack(fill="x", padx=pad)
        tk.Entry(frame_output, textvariable=self.output_pdf, state="readonly").pack(side="left", fill="x", expand=True)
        tk.Button(frame_output, text="Browse", command=self.browse_output_pdf).pack(side="left", padx=5)

        # Generate button
        tk.Button(self, text="Generate PDF", command=self.generate_pdf, bg="#4CAF50", fg="white", font=("Arial", 14)).pack(pady=20)

    def browse_pptx(self):
        file_path = filedialog.askopenfilename(
            title="Choose a PowerPoint file",
            filetypes=[("PowerPoint files", "*.pptx")]
        )
        if file_path:
            self.pptx_path.set(file_path)
            # Auto-suggest output PDF name on pptx select
            base_name = os.path.splitext(os.path.basename(file_path))[0]
            random_suffix = ''.join(random.choices(string.digits, k=4))
            self.output_pdf.set(f"{base_name}_{random_suffix}.pdf")

    def choose_font(self):
        path = choose_font_dialog(self, self.installed_fonts)
        if path:
            self.font_path = path
            font_name = os.path.splitext(os.path.basename(path))[0]
            self.font_name.set(font_name)

    def choose_color(self):
        color_code = colorchooser.askcolor(initialcolor=self._color_to_hex(self.font_color), title="Choose Font Color")
        if color_code[0]:
            self.font_color = tuple(map(int, color_code[0]))
            self.color_display.config(bg=self._color_to_hex(self.font_color))

    def browse_output_pdf(self):
        file_path = filedialog.asksaveasfilename(
            title="Save PDF As",
            defaultextension=".pdf",
            filetypes=[("PDF files", "*.pdf")],
            initialfile=self.output_pdf.get() or "notes.pdf"
        )
        if file_path:
            self.output_pdf.set(file_path)

    def _color_to_hex(self, rgb):
        return "#%02x%02x%02x" % rgb

    def generate_pdf(self):
        if not self.pptx_path.get():
            messagebox.showerror("Error", "Please select a PowerPoint file.")
            return
        if not self.font_path:
            messagebox.showerror("Error", "Please select a font.")
            return
        if not self.output_pdf.get():
            messagebox.showerror("Error", "Please select an output PDF file.")
            return

        try:
            notes = extract_text_from_pptx(self.pptx_path.get())
            margins = {
                'left': self.margin_left.get(),
                'right': self.margin_right.get(),
                'top': self.margin_top.get(),
                'bottom': self.margin_bottom.get()
            }
            font_name = os.path.splitext(os.path.basename(self.font_path))[0]
            write_notes_to_pdf(
                notes,
                self.output_pdf.get(),  # output_path
                self.font_path,         # font_path (was font_name before)
                font_name,              # font_name
                self.font_size.get(),
                margins,
                self.line_spacing.get(),
                self._color_to_hex(self.font_color)  # convert RGB tuple to hex string
            )
            messagebox.showinfo("Success", f"PDF generated at:\n{self.output_pdf.get()}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to generate PDF:\n{e}")

if __name__ == "__main__":
    app = NoteGenUI()
    app.mainloop()
